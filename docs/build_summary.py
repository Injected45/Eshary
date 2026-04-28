"""
Eshary project summary — emits PDF, DOCX, XLSX, PPTX into ./docs/.
All data parsed from project source-of-truth files (pubspec, lib/**/*.dart,
supabase/migrations/*.sql, AndroidManifest.xml). Data lives at the top of
this file as plain Python literals so future updates are mechanical.
"""

from pathlib import Path

# ---------------------------------------------------------------------------
# Parsed data
# ---------------------------------------------------------------------------

OVERVIEW = {
    "name": "Eshary — شركة الرحالة",
    "purpose": (
        "Arabic / RTL financial-operations app for outgoing USD transfers "
        "and currency-buy purchases. Flutter + Supabase port of the original "
        "single-file project_web.html prototype."
    ),
    "primary_language": "Arabic (RTL primary), English secondary",
    "default_locale": "ar",
    "platforms": "Android (com.rahala.eshary) + Flutter Web. Desktop/iOS scaffolds present but not actively targeted.",
    "version": "0.1.0+1",
    "android_label": "شركة الرحالة",
    "package": "com.rahala.eshary",
}

TECH_STACK = [
    ("Flutter SDK", ">=3.24.0", "UI framework, Material 3"),
    ("Dart SDK", ">=3.5.0 <4.0.0", "Language"),
    ("flutter_riverpod", "^2.5.1", "State management — plain providers, no codegen"),
    ("supabase_flutter", "^2.5.6", "Postgres + auth + realtime + session persistence"),
    ("go_router", "^14.2.7", "Declarative routing with auth-state refresh"),
    ("intl", "^0.19.0", "Number / date formatters (ar + en_US)"),
    ("share_plus", "^10.0.2", "Native share sheet"),
    ("pdf", "^3.11.1", "PDF document builder"),
    ("printing", "^5.13.2", "Share / print PDFs"),
    ("shared_preferences", "^2.3.2", "Disk cache + onboarding flag"),
    ("font_awesome_flutter", "^10.7.0", "Icon set"),
    ("google_sign_in", "^6.2.1", "Google ID-token sign-in"),
    ("audioplayers", "^6.1.0", "Asset audio playback"),
    ("flutter_lints", "^4.0.0", "Lints — dev_dependency"),
    ("flutter_localizations", "sdk", "ar + en localization delegates"),
]

ARCHITECTURE = [
    "Feature-folder layout under lib/features/<feature>/{data,domain,presentation}.",
    "core/ holds env, supabase client provider, GoRouter, theme tokens.",
    "shared/ holds glass widgets, formatters, share, JsonCache, PDF export, audio feedback, animated liquid background.",
    "Repositories live in data/, expose async CRUD against Supabase REST + RPC.",
    "Domain models are plain immutable classes with fromJson factories (no freezed).",
    "Riverpod providers expose lists / actions; FutureProviders for queries; Provider<Future<int> Function()> for write actions.",
    "Routing uses GoRouter with refreshListenable bridging supabase auth stream → router redirect.",
    "RTL enforced globally via MaterialApp.builder Directionality wrap.",
    "Read-side offline cache: list responses written into SharedPreferences keyed by user-id, served on network failure.",
    "Atomic write paths use Supabase RPC functions so insert + balance update happen in one Postgres tx.",
]

FEATURES = [
    ("auth", "Email/password + Google sign-in (idToken flow on Android, OAuth redirect on Web). ensureProfile() upsert as belt-and-suspenders."),
    ("onboarding", "3-page intro carousel; onboarding_completed_v1 flag in SharedPreferences."),
    ("splash", "Animated logo splash; routes to onboarding / home / sign-in based on flag + auth session."),
    ("home", "Bottom-nav shell with 4 tabs: خروج / شراء / الأرشيف / الإعدادات. Glass AppBar with dynamic title."),
    ("profile", "Profile card (avatar + email + plan badge), change-password form, sign-out with confirm dialog."),
    ("settings", "Hub with 4 rows: الملف الشخصي / حساباتي / شركات الصرافة / العملاء."),
    ("companies", "User accounts ('حسابات') with one nested exchange each. Country picker + exchange-company dropdown sourced from per-user lists. AccountDetailsScreen shows account + exchange company info."),
    ("exchange_companies", "CRUD list of reusable exchange-office company names per user."),
    ("countries", "Per-user country list, surfaced inline as a glass picker dialog inside AddCompanyDialog."),
    ("clients", "Counterparty CRUD for the currency-buy flow."),
    ("transfers", "Outgoing USD transfer composer: 3 collapsible glass sections (الجهة المنفذة / جهة الاستلام / قيمة التحويل) + saved beneficiaries dialog + daily-log table with click-to-detail and PDF export."),
    ("currency_buy", "USD purchases from clients with auto LYD computation, pending-state and daily-state RPCs, PDF export, archival blocked while pending rows exist."),
    ("archive", "Per-period archive — sold + bought tables, totals, PDF export per section."),
]

SCREENS = [
    ("SplashScreen", "—", "/splash", "lib/features/splash/presentation/splash_screen.dart"),
    ("OnboardingScreen", "إدارة الحوالات / شراء العملات / الأرشيف العام", "/onboarding", "lib/features/onboarding/presentation/onboarding_screen.dart"),
    ("SignInScreen", "مرحباً بعودتك", "/sign-in", "lib/features/auth/presentation/sign_in_screen.dart"),
    ("SignUpScreen", "إنشاء حساب جديد", "/sign-up", "lib/features/auth/presentation/sign_up_screen.dart"),
    ("HomeShell", "<dynamic per tab>", "/", "lib/features/home/presentation/home_shell.dart"),
    ("TransfersScreen", "تنفيذ حوالة جديدة", "/ — tab 0 (خروج)", "lib/features/transfers/presentation/transfers_screen.dart"),
    ("CurrencyBuyScreen", "شراء عملة", "/ — tab 1 (شراء)", "lib/features/currency_buy/presentation/currency_buy_screen.dart"),
    ("ArchiveScreen", "الأرشيف العام", "/ — tab 2 (الأرشيف)", "lib/features/archive/presentation/archive_screen.dart"),
    ("SettingsScreen", "الإعدادات", "/ — tab 3 (الإعدادات)", "lib/features/settings/presentation/settings_screen.dart"),
    ("ProfileDetailsScreen", "الملف الشخصي", "(pushed from Settings)", "lib/features/profile/presentation/profile_details_screen.dart"),
    ("CompaniesScreen", "حساباتي", "(pushed from Settings)", "lib/features/companies/presentation/companies_screen.dart"),
    ("AccountDetailsScreen", "<company.name>", "(pushed from CompaniesScreen)", "lib/features/companies/presentation/account_details_screen.dart"),
    ("ExchangeCompaniesScreen", "شركات الصرافة", "(pushed from Settings)", "lib/features/exchange_companies/presentation/exchange_companies_screen.dart"),
    ("ClientsScreen", "العملاء", "(pushed from Settings)", "lib/features/clients/presentation/clients_screen.dart"),
]

DB_TABLES = [
    ("profiles", "0001", [
        ("id", "uuid", "PK, FK→auth.users(id) ON DELETE CASCADE"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: id = auth.uid()"),
    ("companies", "0001", [
        ("id", "uuid", "PK DEFAULT gen_random_uuid()"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("start_ref", "text", "NOT NULL — reference seed for transfers"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: owner_id = auth.uid()"),
    ("exchanges", "0001 + 0008", [
        ("id", "uuid", "PK"),
        ("company_id", "uuid", "NOT NULL FK→companies(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("balance", "numeric(14,2)", "NOT NULL DEFAULT 0"),
        ("our_code", "text", "nullable"),
        ("country", "text", "nullable — added in 0008"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Derived via companies.owner_id (no direct owner_id column — single source of truth)"),
    ("clients", "0001", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("company", "text", "nullable"),
        ("code", "text", "nullable"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: owner_id = auth.uid()"),
    ("transfers", "0001", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("company_id", "uuid", "NOT NULL FK→companies(id) RESTRICT"),
        ("exchange_id", "uuid", "NOT NULL FK→exchanges(id) RESTRICT"),
        ("beneficiary_name", "text", "NOT NULL"),
        ("beneficiary_account_company", "text", "nullable"),
        ("beneficiary_code", "text", "nullable"),
        ("amount", "numeric(14,2)", "NOT NULL CHECK (amount > 0)"),
        ("reference", "text", "NOT NULL"),
        ("status", "transfer_status", "NOT NULL DEFAULT 'daily' (daily|archived)"),
        ("created_at", "timestamptz", "DEFAULT now()"),
        ("archived_at", "timestamptz", "nullable"),
    ], "Self-only: owner_id = auth.uid()"),
    ("currency_buys", "0001", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("my_company_id", "uuid", "NOT NULL FK→companies(id) RESTRICT"),
        ("exchange_id", "uuid", "NOT NULL FK→exchanges(id) RESTRICT"),
        ("client_id", "uuid", "FK→clients(id) SET NULL"),
        ("client_from_account", "text", "nullable"),
        ("usd_amount", "numeric(14,2)", "NOT NULL CHECK (>0)"),
        ("rate", "numeric(10,4)", "NOT NULL CHECK (>0)"),
        ("lyd_amount", "numeric(14,2)", "NOT NULL CHECK (>=0)"),
        ("status", "currency_buy_status", "NOT NULL DEFAULT 'daily' (pending|daily|archived)"),
        ("created_at", "timestamptz", "DEFAULT now()"),
        ("archived_at", "timestamptz", "nullable"),
    ], "Self-only: owner_id = auth.uid()"),
    ("beneficiaries", "0007", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("account", "text", "nullable"),
        ("code", "text", "nullable"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: owner_id = auth.uid()"),
    ("exchange_companies", "0009", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: ec_*_own — owner_id = auth.uid()"),
    ("countries", "0010", [
        ("id", "uuid", "PK"),
        ("owner_id", "uuid", "NOT NULL FK→profiles(id) CASCADE"),
        ("name", "text", "NOT NULL"),
        ("created_at", "timestamptz", "DEFAULT now()"),
    ], "Self-only: countries_*_own — owner_id = auth.uid()"),
]

RPC_FUNCTIONS = [
    ("archive_daily_transfers(p_owner)", "0003", "Flips owner's daily transfers to archived, stamps archived_at. Returns row count."),
    ("archive_daily_buys(p_owner)", "0003", "Flips owner's daily buys to archived. Refuses if any pending rows exist (raises check_violation)."),
    ("next_reference(p_company_id)", "0003 + 0006", "Returns next transfer reference. 0006 reformat: prefix from start_ref leading non-digits + zero-padded (seed + count) digits. Counts ALL owner transfers, not per-company."),
    ("record_transfer(...)", "0004", "Atomic insert into transfers + decrement exchange.balance. Validates company ownership and exchange/company linkage."),
    ("record_currency_buy(...)", "0004", "Atomic insert into currency_buys (status=daily) + increment exchange.balance by usd_amount."),
    ("record_pending_buy(...)", "0004", "Insert pending currency_buy. No balance mutation."),
    ("handle_new_user()", "0005", "AFTER INSERT trigger on auth.users — autocreates profiles row. SECURITY DEFINER."),
]

DESIGN_TOKENS = [
    ("bgDeep", "#020617", "Background base — slate-950, OLED-friendly"),
    ("bgPanel", "#0B1220", "Stacking surface tone — slate-900-ish"),
    ("glassFill", "#FFFFFF @ 8%", "Glass surface fill"),
    ("glassFillStrong", "#FFFFFF @ 12%", "Stronger glass fill"),
    ("glassBorder", "#FFFFFF @ 12%", "Hairline border"),
    ("glassBorderStrong", "#FFFFFF @ 20%", "Stronger border"),
    ("accent", "#22D3EE", "cyan-400 — primary brand glow"),
    ("accentDim", "#0E7490", "cyan-700"),
    ("positive", "#22C55E", "green-500 — currency buy / income"),
    ("negative", "#EF4444", "red-500 — transfers / outflow"),
    ("warning", "#F59E0B", "amber-500 — pending"),
    ("textHigh", "#F8FAFC", "slate-50"),
    ("textMid", "#CBD5E1", "slate-300"),
    ("textLow", "#94A3B8", "slate-400"),
    ("textDim", "#64748B", "slate-500"),
]

DESIGN_TYPOGRAPHY = [
    ("Almarai", "300 / 400 / 700 / 800", "Primary UI font (default in ThemeData)"),
    ("NotoArabic (Noto Naskh Arabic)", "400 / 500 / 600 / 700", "Alternate Arabic display font, also bundled for PDF export"),
]

DESIGN_WIDGETS = [
    ("GlassCard", "Frosted-glass surface — 18px backdrop blur, translucent fill, hairline border, 24px shadow"),
    ("GlassPanel", "Subtle flatter glass section — 12px blur, 14px radius"),
    ("GlassAppBar", "AppBar with backdrop blur 18 + bgDeep @ 35% alpha"),
    ("showGlassDialog<T>()", "Generalized dialog with slide(0.12→0) + scale(0.94→1) + fade entrance over 280ms easeOutCubic"),
    ("LiquidBackground", "Animated 18s loop. Three radial blobs (accent + positive + violet) blurred 80px on a slate-950 base"),
    ("_CollapsibleSection", "Tap-to-expand glass section with chevron + 220ms easeOutCubic AnimatedSize (transfers screen)"),
]

PROVIDERS = [
    ("supabaseClientProvider", "Provider<SupabaseClient>", "core/supabase_provider.dart"),
    ("authStateChangesProvider", "StreamProvider<AuthState>", "core/supabase_provider.dart"),
    ("currentSessionProvider", "Provider<Session?>", "core/supabase_provider.dart"),
    ("currentUserIdProvider", "Provider<String?>", "core/supabase_provider.dart"),
    ("routerProvider", "Provider<GoRouter>", "core/router.dart"),
    ("authRepositoryProvider", "Provider<AuthRepository>", "features/auth/data"),
    ("onboardingStorageProvider", "Provider<OnboardingStorage>", "features/onboarding/data"),
    ("companiesRepositoryProvider", "Provider", "features/companies/data"),
    ("exchangesRepositoryProvider", "Provider", "features/companies/data"),
    ("companiesListProvider", "FutureProvider<List<Company>>", "features/companies/presentation"),
    ("allExchangesProvider", "FutureProvider<List<Exchange>>", "features/companies/presentation"),
    ("exchangesByCompanyProvider", "FutureProvider.family", "features/companies/presentation"),
    ("nextReferenceProvider", "FutureProvider.family<String,String>", "features/companies/presentation"),
    ("exchangeCompaniesListProvider", "FutureProvider", "features/exchange_companies"),
    ("countriesListProvider", "FutureProvider", "features/countries"),
    ("clientsListProvider", "FutureProvider", "features/clients"),
    ("beneficiariesListProvider", "FutureProvider", "features/transfers"),
    ("dailyTransfersProvider", "FutureProvider", "features/transfers"),
    ("archivedTransfersProvider", "FutureProvider", "features/transfers"),
    ("archiveTransfersActionProvider", "Provider<Future<int> Function()>", "features/transfers"),
    ("dailyBuysProvider", "FutureProvider", "features/currency_buy"),
    ("pendingBuysProvider", "FutureProvider", "features/currency_buy"),
    ("archivedBuysProvider", "FutureProvider", "features/currency_buy"),
    ("archiveBuysActionProvider", "Provider<Future<int> Function()>", "features/currency_buy"),
    ("archivedSoldTotalProvider", "FutureProvider<double>", "features/archive"),
    ("archivedBoughtTotalProvider", "FutureProvider<double>", "features/archive"),
    ("sharedPreferencesProvider", "Provider (overridden in main.dart)", "shared/cache.dart"),
    ("jsonCacheProvider", "Provider<JsonCache>", "shared/cache.dart"),
]

ASSETS = [
    ("assets/fonts/", "Almarai-Light/Regular/Bold/ExtraBold.ttf + NotoNaskhArabic-Regular/Medium/SemiBold/Bold.ttf — registered in pubspec; Almarai-Regular/Bold also loaded by PDF export."),
    ("assets/images/", "background.jpeg — used by sign-in/sign-up screens. google-g.png — Google sign-in icon."),
    ("assets/sounds/", "alert.mp3 — short ding played by playAlert() after save and archive."),
]

RECENT_CHANGES = "—  (git log: branch 'main' has no commits yet)"

# ---------------------------------------------------------------------------
# Output paths
# ---------------------------------------------------------------------------

OUT_DIR = Path(__file__).resolve().parent
PDF_PATH = OUT_DIR / "Eshary_Project_Summary.pdf"
DOCX_PATH = OUT_DIR / "Eshary_Project_Summary.docx"
XLSX_PATH = OUT_DIR / "Eshary_Project_Summary.xlsx"
PPTX_PATH = OUT_DIR / "Eshary_Project_Summary.pptx"

SECTIONS_ORDER = [
    "Overview",
    "Tech Stack",
    "Architecture",
    "Features",
    "Screens",
    "Database",
    "Design System",
    "Recent Changes",
]


# ---------------------------------------------------------------------------
# PDF (reportlab)
# ---------------------------------------------------------------------------

def build_pdf():
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    )

    # Register Almarai so Arabic glyphs render (note: reportlab does not
    # reshape/RTL-reorder; Arabic letters appear visually disconnected but
    # the content is still preserved verbatim).
    fonts_dir = OUT_DIR.parent / "assets" / "fonts"
    pdfmetrics.registerFont(TTFont("Almarai", str(fonts_dir / "Almarai-Regular.ttf")))
    pdfmetrics.registerFont(TTFont("Almarai-Bold", str(fonts_dir / "Almarai-Bold.ttf")))

    styles = getSampleStyleSheet()
    body = ParagraphStyle("body", parent=styles["BodyText"], fontName="Almarai", fontSize=10, leading=14)
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontName="Almarai-Bold", fontSize=18, spaceAfter=10, textColor=colors.HexColor("#0E7490"))
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontName="Almarai-Bold", fontSize=13, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#22D3EE"))
    bullet = ParagraphStyle("bullet", parent=body, leftIndent=14, bulletIndent=2)

    doc = SimpleDocTemplate(str(PDF_PATH), pagesize=A4, leftMargin=1.5 * cm, rightMargin=1.5 * cm, topMargin=1.5 * cm, bottomMargin=1.5 * cm)
    story = []

    story += [Paragraph("Eshary — Project Summary", h1)]
    story += [Paragraph(OVERVIEW["name"], body), Spacer(1, 4)]
    story += [Paragraph(f"Version {OVERVIEW['version']} · package {OVERVIEW['package']} · default locale {OVERVIEW['default_locale']}", body), Spacer(1, 14)]

    story += [Paragraph("1. Overview", h2)]
    for k, v in OVERVIEW.items():
        story.append(Paragraph(f"<b>{k}:</b> {v}", body))
    story.append(Spacer(1, 10))

    story += [Paragraph("2. Tech Stack", h2)]
    rows = [["Package", "Version", "Purpose"]] + [[a, b, c] for (a, b, c) in TECH_STACK]
    story.append(_pdf_table(rows, [4.5 * cm, 3.5 * cm, 9 * cm]))
    story.append(Spacer(1, 10))

    story += [Paragraph("3. Architecture", h2)]
    for line in ARCHITECTURE:
        story.append(Paragraph("• " + line, bullet))
    story.append(Spacer(1, 10))

    story += [PageBreak(), Paragraph("4. Features", h2)]
    rows = [["Folder", "Purpose"]] + [[a, b] for (a, b) in FEATURES]
    story.append(_pdf_table(rows, [4 * cm, 13 * cm]))
    story.append(Spacer(1, 10))

    story += [PageBreak(), Paragraph("5. Screens", h2)]
    rows = [["Screen", "Title (ar)", "Route / location"]] + [[a, b, c] for (a, b, c, _f) in SCREENS]
    story.append(_pdf_table(rows, [5 * cm, 6 * cm, 6 * cm]))
    story.append(Spacer(1, 10))

    story += [PageBreak(), Paragraph("6. Database — Tables", h2)]
    for name, mig, cols, rls in DB_TABLES:
        story.append(Paragraph(f"<b>{name}</b>  ·  migration {mig}", body))
        sub = [["Column", "Type", "Notes"]] + [[c, t, n] for (c, t, n) in cols]
        story.append(_pdf_table(sub, [4 * cm, 4 * cm, 9 * cm]))
        story.append(Paragraph(f"<i>RLS:</i> {rls}", body))
        story.append(Spacer(1, 6))
    story.append(Spacer(1, 8))
    story.append(Paragraph("Database — RPC functions", h2))
    rows = [["Function", "Migration", "Behavior"]] + [[a, b, c] for (a, b, c) in RPC_FUNCTIONS]
    story.append(_pdf_table(rows, [5.5 * cm, 2.5 * cm, 9 * cm]))

    story += [PageBreak(), Paragraph("7. Design System", h2)]
    story.append(Paragraph("<b>Colors</b>", body))
    rows = [["Token", "Value", "Usage"]] + [[a, b, c] for (a, b, c) in DESIGN_TOKENS]
    story.append(_pdf_table(rows, [4 * cm, 5 * cm, 8 * cm]))
    story.append(Spacer(1, 6))
    story.append(Paragraph("<b>Typography</b>", body))
    rows = [["Family", "Weights", "Use"]] + [[a, b, c] for (a, b, c) in DESIGN_TYPOGRAPHY]
    story.append(_pdf_table(rows, [5 * cm, 4 * cm, 8 * cm]))
    story.append(Spacer(1, 6))
    story.append(Paragraph("<b>Reusable widgets</b>", body))
    rows = [["Widget", "Description"]] + [[a, b] for (a, b) in DESIGN_WIDGETS]
    story.append(_pdf_table(rows, [4.5 * cm, 12.5 * cm]))

    story += [PageBreak(), Paragraph("8. Recent Changes", h2)]
    story.append(Paragraph(RECENT_CHANGES, body))

    doc.build(story)


def _pdf_table(rows, col_widths):
    from reportlab.lib import colors
    from reportlab.platypus import Table, TableStyle
    t = Table(rows, colWidths=col_widths, repeatRows=1)
    t.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), "Almarai"),
        ("FONTNAME", (0, 0), (-1, 0), "Almarai-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0B1220")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#94A3B8")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    return t


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------

def build_docx():
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.shared import Pt, RGBColor, Cm
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10)

    def add_h(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0x0E, 0x74, 0x90)
        return h

    def add_para(text, rtl=False):
        p = doc.add_paragraph(text)
        if rtl:
            pPr = p._p.get_or_add_pPr()
            bidi = OxmlElement("w:bidi")
            bidi.set(qn("w:val"), "1")
            pPr.append(bidi)
            p.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
        return p

    def add_table(headers, rows):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = "Light Grid Accent 1"
        hdr = t.rows[0].cells
        for i, h in enumerate(headers):
            hdr[i].text = h
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                t.rows[r].cells[c].text = str(val)
        doc.add_paragraph()
        return t

    add_h("Eshary — Project Summary", level=0)
    add_para(OVERVIEW["name"], rtl=True)
    add_para(f"Version {OVERVIEW['version']} · package {OVERVIEW['package']} · default locale {OVERVIEW['default_locale']}")

    add_h("1. Overview", level=1)
    for k, v in OVERVIEW.items():
        p = doc.add_paragraph()
        p.add_run(f"{k}: ").bold = True
        p.add_run(str(v))

    add_h("2. Tech Stack", level=1)
    add_table(["Package", "Version", "Purpose"], TECH_STACK)

    add_h("3. Architecture", level=1)
    for line in ARCHITECTURE:
        doc.add_paragraph(line, style="List Bullet")

    add_h("4. Features", level=1)
    add_table(["Folder", "Purpose"], FEATURES)

    add_h("5. Screens", level=1)
    add_table(["Screen", "Arabic title", "Route / location", "File"],
              [(a, b, c, f) for (a, b, c, f) in SCREENS])

    add_h("6. Database — Tables", level=1)
    for name, mig, cols, rls in DB_TABLES:
        h = doc.add_paragraph()
        r = h.add_run(f"{name} ")
        r.bold = True
        h.add_run(f"(migration {mig})")
        add_table(["Column", "Type", "Notes"], cols)
        rls_p = doc.add_paragraph()
        rls_p.add_run("RLS: ").italic = True
        rls_p.add_run(rls)

    add_h("Database — RPC functions", level=2)
    add_table(["Function", "Migration", "Behavior"], RPC_FUNCTIONS)

    add_h("7. Design System", level=1)
    doc.add_paragraph("Colors", style="Heading 3")
    add_table(["Token", "Value", "Usage"], DESIGN_TOKENS)
    doc.add_paragraph("Typography", style="Heading 3")
    add_table(["Family", "Weights", "Use"], DESIGN_TYPOGRAPHY)
    doc.add_paragraph("Reusable widgets", style="Heading 3")
    add_table(["Widget", "Description"], DESIGN_WIDGETS)
    doc.add_paragraph("Riverpod providers", style="Heading 3")
    add_table(["Provider", "Type", "Location"], PROVIDERS)
    doc.add_paragraph("Assets", style="Heading 3")
    add_table(["Path", "Contents"], ASSETS)

    add_h("8. Recent Changes", level=1)
    doc.add_paragraph(RECENT_CHANGES)

    doc.save(str(DOCX_PATH))


# ---------------------------------------------------------------------------
# XLSX (openpyxl)
# ---------------------------------------------------------------------------

def build_xlsx():
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill, Border, Side

    wb = Workbook()
    wb.remove(wb.active)

    head_fill = PatternFill("solid", fgColor="0B1220")
    head_font = Font(bold=True, color="FFFFFF")
    border = Border(*[Side(style="thin", color="94A3B8")] * 4)

    def sheet(name, headers, rows):
        ws = wb.create_sheet(title=name)
        ws.append(headers)
        for c in ws[1]:
            c.fill = head_fill
            c.font = head_font
            c.alignment = Alignment(vertical="center", horizontal="center")
            c.border = border
        for row in rows:
            ws.append(list(row))
        for r_idx in range(2, ws.max_row + 1):
            for c in ws[r_idx]:
                c.alignment = Alignment(vertical="top", wrap_text=True)
                c.border = border
        for col in ws.columns:
            letter = col[0].column_letter
            max_len = max((len(str(c.value)) if c.value is not None else 0) for c in col)
            ws.column_dimensions[letter].width = min(max(14, max_len + 2), 60)
        ws.freeze_panes = "A2"

    sheet("Overview", ["Field", "Value"], list(OVERVIEW.items()))
    sheet("TechStack", ["Package", "Version", "Purpose"], TECH_STACK)
    sheet("Architecture", ["Note"], [(line,) for line in ARCHITECTURE])
    sheet("Features", ["Folder", "Purpose"], FEATURES)
    sheet("Screens", ["Screen", "Arabic title", "Route / location", "File"], SCREENS)

    db_rows = []
    for name, mig, cols, rls in DB_TABLES:
        for col, ty, notes in cols:
            db_rows.append((name, mig, col, ty, notes, rls))
    sheet("DBTables", ["Table", "Migration", "Column", "Type", "Notes", "RLS"], db_rows)
    sheet("RPCFunctions", ["Function", "Migration", "Behavior"], RPC_FUNCTIONS)

    rls_rows = [(t[0], t[3]) for t in DB_TABLES]
    sheet("RLSPolicies", ["Table", "Policy summary"], rls_rows)

    sheet("DesignTokens", ["Token", "Value", "Usage"], DESIGN_TOKENS)
    sheet("Typography", ["Family", "Weights", "Use"], DESIGN_TYPOGRAPHY)
    sheet("Widgets", ["Widget", "Description"], DESIGN_WIDGETS)
    sheet("Providers", ["Provider", "Type", "Location"], PROVIDERS)
    sheet("Assets", ["Path", "Contents"], ASSETS)
    sheet("RecentChanges", ["Source", "Output"], [("git log --oneline -n 20", RECENT_CHANGES)])

    wb.save(str(XLSX_PATH))


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(0x02, 0x06, 0x17)
    ACCENT = RGBColor(0x22, 0xD3, 0xEE)
    TEXT_HIGH = RGBColor(0xF8, 0xFA, 0xFC)
    TEXT_LOW = RGBColor(0x94, 0xA3, 0xB8)

    def blank_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        return s

    def title_block(slide, title, subtitle=None):
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(16)
            r2.font.color.rgb = TEXT_LOW

    def bullets(slide, items, top=Inches(1.6), font_size=14, max_per_box=18):
        if not items:
            return
        chunks = [items[i:i + max_per_box] for i in range(0, len(items), max_per_box)]
        n = len(chunks)
        widths = [Inches(12 / n) for _ in chunks]
        x = Inches(0.6)
        for chunk, w in zip(chunks, widths):
            tx = slide.shapes.add_textbox(x, top, w, Inches(5.5))
            tf = tx.text_frame
            tf.word_wrap = True
            for i, item in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = "• " + str(item)
                r.font.size = Pt(font_size)
                r.font.color.rgb = TEXT_HIGH
                p.space_after = Pt(4)
            x += w + Inches(0.1)

    def kv_block(slide, pairs, top=Inches(1.6)):
        tx = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(5.5))
        tf = tx.text_frame
        tf.word_wrap = True
        first = True
        for k, v in pairs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = f"{k}: "
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = ACCENT
            r2 = p.add_run()
            r2.text = str(v)
            r2.font.size = Pt(14)
            r2.font.color.rgb = TEXT_HIGH

    def table_block(slide, headers, rows, top=Inches(1.6), col_widths=None):
        cols = len(headers)
        rcount = len(rows) + 1
        width = Inches(12)
        height = Inches(5.5)
        tbl_shape = slide.shapes.add_table(rcount, cols, Inches(0.6), top, width, height)
        tbl = tbl_shape.table
        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = w
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(11)
                    r.font.color.rgb = TEXT_HIGH
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x20)
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)
                        r.font.color.rgb = TEXT_HIGH

    # ---- Slide 1 — title
    s = blank_slide()
    title_block(s, OVERVIEW["name"], "Eshary — Project Summary")
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = OVERVIEW["purpose"]
    r.font.size = Pt(18)
    r.font.color.rgb = TEXT_HIGH
    p2 = tf.add_paragraph()
    p2.space_before = Pt(20)
    r2 = p2.add_run()
    r2.text = (
        f"Version {OVERVIEW['version']}    ·    package {OVERVIEW['package']}    "
        f"·    default locale {OVERVIEW['default_locale']}"
    )
    r2.font.size = Pt(13)
    r2.font.color.rgb = TEXT_LOW

    # ---- Slide 2 — Overview
    s = blank_slide()
    title_block(s, "1. Overview")
    kv_block(s, list(OVERVIEW.items()))

    # ---- Slide 3 — Tech Stack
    s = blank_slide()
    title_block(s, "2. Tech Stack")
    table_block(s, ["Package", "Version", "Purpose"], TECH_STACK,
                col_widths=[Inches(3.5), Inches(2.0), Inches(6.5)])

    # ---- Slide 4 — Architecture
    s = blank_slide()
    title_block(s, "3. Architecture")
    bullets(s, ARCHITECTURE)

    # ---- Slide 5 — Features
    s = blank_slide()
    title_block(s, "4. Features")
    table_block(s, ["Folder", "Purpose"], FEATURES,
                col_widths=[Inches(2.5), Inches(9.5)])

    # ---- Slide 6 — Screens
    s = blank_slide()
    title_block(s, "5. Screens")
    table_block(s, ["Screen", "Arabic title", "Route / location"],
                [(a, b, c) for (a, b, c, _f) in SCREENS],
                col_widths=[Inches(3.5), Inches(3.5), Inches(5.0)])

    # ---- Slide 7 — DB Tables (overview)
    s = blank_slide()
    title_block(s, "6. Database — Tables")
    rows = [(name, mig, len(cols), rls) for (name, mig, cols, rls) in DB_TABLES]
    table_block(s, ["Table", "Migration", "Cols", "RLS summary"], rows,
                col_widths=[Inches(2.5), Inches(2.0), Inches(1.0), Inches(6.5)])

    # ---- Slide 8 — RPC functions
    s = blank_slide()
    title_block(s, "Database — RPC functions")
    table_block(s, ["Function", "Migration", "Behavior"], RPC_FUNCTIONS,
                col_widths=[Inches(3.5), Inches(1.5), Inches(7.0)])

    # ---- Slide 9 — Design tokens
    s = blank_slide()
    title_block(s, "7. Design System — Colors")
    table_block(s, ["Token", "Value", "Usage"], DESIGN_TOKENS,
                col_widths=[Inches(3.0), Inches(3.0), Inches(6.0)])

    # ---- Slide 10 — Typography + Widgets
    s = blank_slide()
    title_block(s, "Design System — Typography & Widgets")
    table_block(s, ["Family / Widget", "Detail"],
                [(a, b + " · " + c) for (a, b, c) in DESIGN_TYPOGRAPHY] + DESIGN_WIDGETS,
                col_widths=[Inches(3.5), Inches(8.5)])

    # ---- Slide 11 — Providers
    s = blank_slide()
    title_block(s, "State — Riverpod Providers")
    table_block(s, ["Provider", "Type", "Location"], PROVIDERS,
                col_widths=[Inches(4.0), Inches(3.5), Inches(4.5)])

    # ---- Slide 12 — Assets
    s = blank_slide()
    title_block(s, "Assets")
    table_block(s, ["Path", "Contents"], ASSETS,
                col_widths=[Inches(2.5), Inches(9.5)])

    # ---- Slide 13 — Recent Changes
    s = blank_slide()
    title_block(s, "8. Recent Changes")
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = RECENT_CHANGES
    r.font.size = Pt(18)
    r.font.color.rgb = TEXT_HIGH

    prs.save(str(PPTX_PATH))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    build_pdf()
    print(f"PDF   OK{PDF_PATH.name}  ({PDF_PATH.stat().st_size:,} bytes)")
    build_docx()
    print(f"DOCX  OK{DOCX_PATH.name}  ({DOCX_PATH.stat().st_size:,} bytes)")
    build_xlsx()
    print(f"XLSX  OK{XLSX_PATH.name}  ({XLSX_PATH.stat().st_size:,} bytes)")
    build_pptx()
    print(f"PPTX  OK{PPTX_PATH.name}  ({PPTX_PATH.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
